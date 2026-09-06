"""Unit tests for Phase 4.1 structured PDF/PPTX intake.

PDF and PPTX fixtures are generated in-test (a hand-built two-page PDF and
a python-pptx slide deck), so no binary files enter the repository.  Parser
behaviour is exercised without a database; persistence, jobs, and the worker
step run against PostgreSQL when ``V2_TEST_DATABASE_URL`` is set.
"""

from __future__ import annotations

import os
import struct
import tempfile
import unittest
import zlib
from io import BytesIO

import psycopg
from psycopg.rows import dict_row

from v2.documents import (
    DocumentConflict,
    DocumentError,
    claim_document_job,
    create_version,
    detect_file_type,
    get_blocks,
    get_document_job,
    get_version,
    latest_document_job,
    list_versions,
    parse_file,
    retry_document_job,
    save_parsed_blocks,
    storage_dir,
    version_file_path,
)


DATABASE_URL = os.getenv("V2_TEST_DATABASE_URL", "").strip()


# -- in-test fixtures --------------------------------------------------------


def _png_bytes() -> bytes:
    signature = b"\x89PNG\r\n\x1a\n"

    def chunk(kind: bytes, data: bytes) -> bytes:
        body = struct.pack(">I", len(data)) + kind + data
        return body + struct.pack(">I", zlib.crc32(kind + data) & 0xFFFFFFFF)

    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    raw = b"\x00\xff\x00\x00"
    return (
        signature
        + chunk(b"IHDR", ihdr)
        + chunk(b"IDAT", zlib.compress(raw))
        + chunk(b"IEND", b"")
    )


def _pdf_bytes() -> bytes:
    """Two pages: a titled text page, then an image-only (scanned-like) page."""

    page1 = (
        b"BT /F2 24 Tf 72 720 Td (Quick Start Guide) Tj ET\n"
        b"BT /F1 12 Tf 72 690 Td (Add a user before anything else.) Tj ET\n"
        b"BT /F1 12 Tf 72 672 Td (Open Users and tap Add to create one.) Tj ET\n"
        b"BT /F2 16 Tf 72 640 Td (Fingerprints) Tj ET\n"
        b"BT /F1 12 Tf 72 622 Td (Enroll two fingerprints per user.) Tj ET\n"
    )
    page2 = b"q 100 0 0 100 72 700 cm /Im1 Do Q\n"
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R 4 0 R] /Count 2 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 5 0 R /F2 6 0 R >> >> /Contents 7 0 R >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /XObject << /Im1 8 0 R >> >> /Contents 9 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica-Bold >>",
        b"<< /Length %d >>\nstream\n" % len(page1) + page1 + b"endstream",
        b"<< /Type /XObject /Subtype /Image /Width 1 /Height 1 "
        b"/ColorSpace /DeviceGray /BitsPerComponent 8 /Length 1 >>\nstream\n\x00\nendstream",
        b"<< /Length %d >>\nstream\n" % len(page2) + page2 + b"endstream",
    ]
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for number, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += b"%d 0 obj\n" % number + body + b"\nendobj\n"
    xref_at = len(out)
    out += b"xref\n0 %d\n" % (len(objects) + 1)
    out += b"0000000000 65535 f \n"
    for offset in offsets:
        out += b"%010d 00000 n \n" % offset
    out += (
        b"trailer\n<< /Size %d /Root 1 0 R >>\nstartxref\n%d\n%%%%EOF\n"
        % (len(objects) + 1, xref_at)
    )
    return bytes(out)


def _pptx_bytes() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    # Slide 1: title placeholder.
    slide1 = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide1.shapes.title.text = "Night Imaging Training"
    # Slide 2: title + textbox + merged-header table + speaker notes.
    slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide2.shapes.title.text = "Models and Modes"
    box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    box.text_frame.text = "Use infrared when no extra light is available."
    graphic = slide2.shapes.add_table(3, 2, Inches(0.5), Inches(2.8), Inches(9), Inches(2))
    table = graphic.table
    table.cell(0, 0).text = "Imaging Mode"
    table.cell(0, 0).merge(table.cell(0, 1))
    table.cell(1, 0).text = "Color"
    table.cell(1, 1).text = "needs light"
    table.cell(2, 0).text = "Infrared"
    table.cell(2, 1).text = "night default"
    slide2.notes_slide.notes_text_frame.text = "Remind engineers about light first."
    # Slide 3: picture only.
    slide3 = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide3.shapes.add_picture(BytesIO(_png_bytes()), Inches(1), Inches(1))
    # Slide 4: intentionally empty.
    presentation.slides.add_slide(presentation.slide_layouts[6])
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _write_temp(suffix: str, data: bytes) -> str:
    handle, path = tempfile.mkstemp(suffix=suffix)
    with os.fdopen(handle, "wb") as stream:
        stream.write(data)
    return path


# -- parser unit tests (no database) -----------------------------------------


class PdfParserTest(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.path = _write_temp(".pdf", _pdf_bytes())

    @classmethod
    def tearDownClass(cls):
        os.unlink(cls.path)

    def test_text_page_structure(self):
        blocks = parse_file(self.path, "pdf")
        by_key = {block["block_key"]: block for block in blocks}
        self.assertEqual(len(blocks), len(by_key), "block keys must be unique")
        page1 = [block for block in blocks if block["page_no"] == 1]
        kinds = [block["block_type"] for block in page1]
        self.assertIn("heading", kinds)
        self.assertIn("paragraph", kinds)
        title = next(block for block in page1 if block["text"] == "Quick Start Guide")
        self.assertEqual(title["block_type"], "heading")
        self.assertEqual(title["section_path"], ["Quick Start Guide"])
        sub = next(block for block in page1 if block["text"] == "Fingerprints")
        self.assertEqual(sub["section_path"], ["Quick Start Guide", "Fingerprints"])
        body = next(block for block in page1 if "tap Add" in block["text"])
        self.assertEqual(body["section_path"], ["Quick Start Guide"])
        self.assertTrue(all(block["slide_no"] is None for block in blocks))

    def test_image_only_page_is_flagged_not_understood(self):
        blocks = parse_file(self.path, "pdf")
        page2 = [block for block in blocks if block["page_no"] == 2]
        self.assertTrue(page2, "every page must have a destination")
        images = [block for block in page2 if block["block_type"] == "image"]
        self.assertEqual(len(images), 1)
        (image,) = images
        self.assertEqual(image["text"], "")
        self.assertEqual(image["needs_review_reason"], "image_only_page")
        (asset,) = image["assets"]
        self.assertEqual(asset["interpretation_status"], "unexplained")
        self.assertTrue(asset["asset_ref"].startswith("page-2-img-"))
        self.assertFalse(any(
            block["block_type"] == "paragraph" and block["text"] for block in page2
        ))


class PptxParserTest(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.path = _write_temp(".pptx", _pptx_bytes())

    @classmethod
    def tearDownClass(cls):
        os.unlink(cls.path)

    def test_slide_titles_and_order(self):
        blocks = parse_file(self.path, "pptx")
        by_key = {block["block_key"]: block for block in blocks}
        self.assertEqual(len(blocks), len(by_key), "block keys must be unique")
        self.assertEqual(
            [block["slide_no"] for block in blocks],
            sorted(block["slide_no"] for block in blocks),
        )
        title = next(block for block in blocks if block["block_key"] == "s1-title")
        self.assertEqual(title["text"], "Night Imaging Training")
        self.assertEqual(title["block_type"], "heading")

    def test_table_cells_header_and_merge_span(self):
        blocks = parse_file(self.path, "pptx")
        tables = [block for block in blocks if block["block_type"] == "table"]
        self.assertEqual(len(tables), 1)
        (table,) = tables
        self.assertEqual(table["slide_no"], 2)
        layout = table["layout"]
        self.assertEqual(layout["rows"], 3)
        self.assertEqual(layout["columns"], 2)
        self.assertEqual(layout["header_row"], 0)
        head = next(cell for cell in layout["cells"] if cell["row"] == 0 and cell["col"] == 0)
        self.assertEqual(head["text"], "Imaging Mode")
        self.assertEqual(head["col_span"], 2)
        self.assertIn("Color | needs light", table["text"])
        self.assertIn("Infrared | night default", table["text"])

    def test_notes_and_image_and_empty_slide(self):
        blocks = parse_file(self.path, "pptx")
        notes = [block for block in blocks if block["block_type"] == "note"]
        self.assertEqual(len(notes), 1)
        self.assertEqual(notes[0]["slide_no"], 2)
        self.assertIn("light first", notes[0]["text"])
        images = [block for block in blocks if block["block_type"] == "image"]
        self.assertEqual(len(images), 1)
        self.assertEqual(images[0]["slide_no"], 3)
        (asset,) = images[0]["assets"]
        self.assertEqual(asset["interpretation_status"], "unexplained")
        self.assertTrue(isinstance(asset["image_bytes"], bytes))
        empty = [block for block in blocks if block["slide_no"] == 4]
        self.assertEqual(len(empty), 1)
        self.assertEqual(empty[0]["needs_review_reason"], "empty_slide")


class UploadValidationTest(unittest.TestCase):
    def test_magic_detection_beats_extension(self):
        self.assertEqual(detect_file_type("manual.pdf", _pdf_bytes()), "pdf")
        self.assertEqual(detect_file_type("deck.pptx", _pptx_bytes()), "pptx")
        self.assertEqual(detect_file_type("renamed.bin", _pdf_bytes()), "pdf")
        with self.assertRaises(DocumentError):
            detect_file_type("notes.txt", b"just some text")
        with self.assertRaises(DocumentError):
            parse_file("/nonexistent", "zip")


# -- PostgreSQL integration ---------------------------------------------------


@unittest.skipUnless(DATABASE_URL, "set V2_TEST_DATABASE_URL to run PostgreSQL integration tests")
class V2DocumentsPostgresTest(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.conn = psycopg.connect(DATABASE_URL, row_factory=dict_row)

    def tearDown(self):
        try:
            self.conn.rollback()
            with psycopg.connect(DATABASE_URL, row_factory=dict_row) as conn:
                with conn.cursor() as cur:
                    cur.execute(
                        "DELETE FROM v2_document_jobs WHERE version_id IN ("
                        "SELECT id FROM v2_document_versions WHERE document_key LIKE 'ZZDOC %')"
                    )
                    cur.execute(
                        "DELETE FROM v2_document_blocks WHERE version_id IN ("
                        "SELECT id FROM v2_document_versions WHERE document_key LIKE 'ZZDOC %')"
                    )
                    cur.execute(
                        "DELETE FROM v2_raw_evidence WHERE source_label LIKE 'ZZDOC %'"
                    )
                    cur.execute(
                        "DELETE FROM v2_document_versions WHERE document_key LIKE 'ZZDOC %'"
                    )
        finally:
            self.conn.rollback()
            self.conn.close()
            self.tmp.cleanup()

    def _upload(self, key, label, data, filename="doc.pdf", **kwargs):
        params = {
            "base_dir": self.tmp.name, "document_key": key,
            "version_label": label, "filename": filename, "content": data,
        }
        params.update(kwargs)
        version, created = create_version(self.conn, **params)
        self.conn.commit()
        return version, created

    def test_reupload_is_idempotent_and_label_takeover_conflicts(self):
        data = _pdf_bytes()
        first, created = self._upload("ZZDOC manual", "v1", data, filename="m.pdf")
        self.assertTrue(created)
        self.assertEqual(first["file_type"], "pdf")
        self.assertEqual(first["status"], "uploaded")
        again, created = self._upload("ZZDOC manual", "v1", data, filename="m.pdf")
        self.assertFalse(created)
        self.assertEqual(int(again["id"]), int(first["id"]))
        jobs = self.conn.execute(
            "SELECT count(*) AS n FROM v2_document_jobs WHERE version_id=%s",
            (int(first["id"]),),
        ).fetchone()["n"]
        self.assertEqual(jobs, 1)
        with self.assertRaises(DocumentConflict):
            self._upload("ZZDOC manual", "v1", _pptx_bytes(), filename="d.pptx")
        self.conn.rollback()
        with self.assertRaises(DocumentError):
            self._upload("ZZDOC bad", "v1", b"nope", filename="x.pdf")
        self.conn.rollback()

    def test_parse_job_processes_blocks_with_evidence(self):
        from v2.document_processing import process_document_job

        version, _ = self._upload("ZZDOC parse", "v1", _pdf_bytes(), filename="m.pdf")
        with psycopg.connect(DATABASE_URL, row_factory=dict_row) as conn:
            job = claim_document_job(conn, ("parse",))
            conn.commit()
        self.assertIsNotNone(job)
        process_document_job(
            int(job["id"]),
            db_factory=lambda: psycopg.connect(DATABASE_URL, row_factory=dict_row),
            base_dir=self.tmp.name,
        )
        with psycopg.connect(DATABASE_URL, row_factory=dict_row) as conn:
            stored = get_version(conn, int(version["id"]))
            self.assertEqual(stored["status"], "parsed")
            blocks = get_blocks(conn, int(version["id"]))
            self.assertGreaterEqual(len(blocks), 6)
            keyed = {block["block_key"]: block for block in blocks}
            self.assertEqual(len(keyed), len(blocks))
            flagged = [block for block in blocks if block["processing_state"] == "needs_review"]
            self.assertTrue(flagged, "image-only page must be flagged")
            texts = [block for block in blocks if block["evidence_text"]]
            self.assertTrue(texts, "text blocks keep raw evidence")
            done = get_document_job(conn, int(job["id"]))
            self.assertEqual(done["status"], "completed")
            self.assertGreaterEqual(int(done["result_summary"]["blocks"]), 6)

    def test_pptx_assets_are_preserved_as_files(self):
        from v2.document_processing import process_document_job

        version, _ = self._upload("ZZDOC assets", "v1", _pptx_bytes(), filename="d.pptx")
        with psycopg.connect(DATABASE_URL, row_factory=dict_row) as conn:
            job = claim_document_job(conn, ("parse",))
            conn.commit()
        process_document_job(
            int(job["id"]),
            db_factory=lambda: psycopg.connect(DATABASE_URL, row_factory=dict_row),
            base_dir=self.tmp.name,
        )
        with psycopg.connect(DATABASE_URL, row_factory=dict_row) as conn:
            blocks = get_blocks(conn, int(version["id"]))
        images = [block for block in blocks if block["block_type"] == "image"]
        self.assertEqual(len(images), 1)
        (asset,) = images[0]["assets"]
        target = storage_dir(self.tmp.name) / "assets" / asset["asset_ref"].split("/", 1)[1]
        self.assertTrue(target.is_file(), "image bytes must be preserved on disk")
        self.assertEqual(asset["interpretation_status"], "unexplained")

    def test_garbage_bytes_fail_cleanly_and_retry(self):
        from v2.document_processing import process_document_job

        version, _ = self._upload(
            "ZZDOC garbage", "v1", b"%PDF-1.4\n%garbage\n", filename="g.pdf",
        )
        with psycopg.connect(DATABASE_URL, row_factory=dict_row) as conn:
            job = claim_document_job(conn, ("parse",))
            conn.commit()
        factory = lambda: psycopg.connect(DATABASE_URL, row_factory=dict_row)
        process_document_job(int(job["id"]), db_factory=factory, base_dir=self.tmp.name)
        with psycopg.connect(DATABASE_URL, row_factory=dict_row) as conn:
            failed = get_document_job(conn, int(job["id"]))
            # A corrupt parse is transient-class: requeued with backoff, and
            # the version is NOT marked parsed.
            self.assertIn(failed["status"], ("queued", "failed"))
            self.assertTrue(failed["error"])
            stored = get_version(conn, int(version["id"]))
            self.assertNotEqual(stored["status"], "parsed")
            if failed["status"] == "failed":
                retried = retry_document_job(conn, int(job["id"]))
                conn.commit()
                self.assertEqual(retried["status"], "queued")
            latest = latest_document_job(conn, int(version["id"]), "parse")
            self.assertEqual(int(latest["id"]), int(job["id"]))

    def test_stored_path_escape_is_unavailable(self):
        from v2.documents import DocumentNotFound

        version, _ = self._upload("ZZDOC escape", "v1", _pdf_bytes(), filename="m.pdf")
        evil = dict(version)
        evil["stored_path"] = "../outside.pdf"
        with self.assertRaises(DocumentNotFound):
            version_file_path(self.tmp.name, evil)
        versions = list_versions(self.conn)
        self.assertTrue(any(item["document_key"].startswith("ZZDOC") for item in versions))


class DummyCursor:
    def __enter__(self):
        return self

    def __exit__(self, *_):
        return False

    def execute(self, query, params=()):
        text = " ".join(str(query).split())
        if "FROM v2_document_jobs WHERE idempotency_key=" in text:
            self._result = [{"id": 77}]
        else:  # pragma: no cover - upload route only runs the lookup above
            raise AssertionError(f"unexpected query: {text[:80]}")

    def fetchone(self):
        return self._result[0] if self._result else None


class DummyConn:
    def __enter__(self):
        return self

    def __exit__(self, *_):
        return False

    def cursor(self):
        return DummyCursor()


def _version_row():
    return {
        "id": 9, "document_key": "manual", "version_label": "v1",
        "sha256": "abc", "file_name": "m.pdf", "file_type": "pdf",
        "file_size": 10, "title": "Manual", "applicability": {},
        "source_authenticity": "unverified", "parser_version": "p",
        "status": "parsed", "block_count": 3,
        "created_at": None, "updated_at": None,
    }


class DocumentApiTest(unittest.TestCase):
    """Route-level tests: HTTP mapping, not service logic."""

    def setUp(self):
        from unittest.mock import patch

        import app as app_module

        self.app_module = app_module
        self._previous_api_key = app_module.settings["api_key"]
        app_module.settings["api_key"] = "test-key"
        self.db_patch = patch.object(app_module, "db", return_value=DummyConn())
        self.db_patch.start()
        self.addCleanup(self.db_patch.stop)
        self.addCleanup(app_module.settings.__setitem__, "api_key", self._previous_api_key)

    def _patch(self, name, value):
        from unittest.mock import patch

        patcher = patch.object(self.app_module, name, value)
        patcher.start()
        self.addCleanup(patcher.stop)

    def _upload_file(self, data: bytes, filename: str):
        from starlette.datastructures import UploadFile

        return UploadFile(file=BytesIO(data), filename=filename)

    def test_upload_shape_and_errors(self):
        import asyncio

        from fastapi import HTTPException

        from app import v2_upload_document
        from v2.documents import DocumentConflict, DocumentError

        version = _version_row()
        self._patch("create_version", lambda conn, **_: (version, True))
        response = asyncio.run(v2_upload_document(
            file=self._upload_file(_pdf_bytes(), "m.pdf"),
            document_key="manual", version_label="v1", title="",
            applicability="", source_authenticity="unverified",
            x_api_key="test-key",
        ))
        self.assertEqual(response["version_id"], 9)
        self.assertTrue(response["created"])

        def bad_magic(conn, **_):
            raise DocumentError("neither PDF nor PPTX")

        self._patch("create_version", bad_magic)
        with self.assertRaises(HTTPException) as caught:
            asyncio.run(v2_upload_document(
                file=self._upload_file(b"nope", "x.pdf"),
                document_key="manual", version_label="v1", title="",
                applicability="", source_authenticity="unverified",
                x_api_key="test-key",
            ))
        self.assertEqual(caught.exception.status_code, 400)

        def taken(conn, **_):
            raise DocumentConflict("label taken")

        self._patch("create_version", taken)
        with self.assertRaises(HTTPException) as caught:
            asyncio.run(v2_upload_document(
                file=self._upload_file(_pdf_bytes(), "m.pdf"),
                document_key="manual", version_label="v1", title="",
                applicability="", source_authenticity="unverified",
                x_api_key="test-key",
            ))
        self.assertEqual(caught.exception.status_code, 409)

    def test_documents_envelope_lists_versions(self):
        from app import v2_documents

        self._patch("list_documents", lambda conn: [])
        self._patch("list_versions", lambda conn: [_version_row()])
        response = v2_documents(x_api_key="test-key")
        self.assertEqual(response["total"], 0)
        self.assertEqual(len(response["versions"]), 1)
        self.assertEqual(response["versions"][0]["version_id"], 9)

    def test_version_blocks_and_job_shapes(self):
        from fastapi import HTTPException

        from app import v2_document_blocks, v2_document_job, v2_retry_document_job

        block = {
            "id": 1, "block_key": "p1-1", "ord": 1, "section_path": [],
            "page_no": 1, "slide_no": None, "block_type": "heading",
            "raw_evidence_id": 2, "evidence_text": "Guide",
            "content_hash": "h", "layout": {}, "assets": [],
            "processing_state": "pending", "state_reason": "",
        }
        self._patch("get_version", lambda conn, _: _version_row())
        self._patch("get_blocks", lambda conn, _: [block])
        self._patch("latest_document_job",
                    lambda conn, *_, **__: {"id": 4, "status": "completed",
                                            "error": "", "result_summary": {"blocks": 1}})
        detail = v2_document_blocks(9, x_api_key="test-key")
        self.assertEqual(detail["total"], 1)
        self.assertEqual(detail["parse_job"]["job_id"], 4)
        self.assertEqual(detail["items"][0]["evidence_text"], "Guide")

        self._patch("get_document_job", lambda conn, _: None)
        with self.assertRaises(HTTPException) as caught:
            v2_document_job(4242, x_api_key="test-key")
        self.assertEqual(caught.exception.status_code, 404)

        self._patch("get_document_job",
                    lambda conn, _: {"id": 4, "version_id": 9, "stage": "parse",
                                     "status": "failed", "attempts": 1,
                                     "result_summary": {}, "error": "boom",
                                     "created_at": None, "updated_at": None})
        job = v2_document_job(4, x_api_key="test-key")
        self.assertEqual(job["status"], "failed")

        self._patch("retry_document_job",
                    lambda conn, _: {"id": 4, "status": "queued"})
        retried = v2_retry_document_job(4, x_api_key="test-key")
        self.assertEqual(retried["status"], "queued")


if __name__ == "__main__":
    unittest.main()
