"""Phase 4.1 structured PDF/PPTX intake: file versions and structure blocks.

PDF and PPTX are both first-class inputs; nobody has to convert anything.
Every upload becomes an immutable ``v2_document_versions`` row (bytes are
never overwritten) plus a queued ``parse`` job.  Parsing extracts the
document structure -- pages/slides, headings, paragraphs, tables with
headers and merge spans, images, and speaker notes -- into
``v2_document_blocks``.  Block text lives in ``v2_raw_evidence`` (the single
original-text store); blocks only point at it.

Honesty rules: a scanned/image-only page is marked ``needs_review``, never
pretended understood.  Charts, SmartArt, and diagrams that cannot be
reliably recovered are stored as unexplained assets, never as product
facts.  Nothing parsed here becomes answerable Knowledge in 4.1; promotion
to Knowledge proposals arrives in Phase 4.2.
"""

from __future__ import annotations

import hashlib
import json
import logging
from pathlib import Path
from typing import Any

from psycopg.types.json import Jsonb

log = logging.getLogger("aihelper.v2.documents")

FILE_TYPES = ("pdf", "pptx")
PDF_MAGIC = b"%PDF"
PPTX_MAGIC = b"PK\x03\x04"
MIME_TYPES = {
    "pdf": "application/pdf",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentationml",
}
# Upload and parse budgets: one request must never wedge the service.
MAX_UPLOAD_BYTES = 50 * 1024 * 1024
MAX_PDF_PAGES = 500
MAX_PPTX_SLIDES = 500
MAX_ASSET_BYTES = 100 * 1024 * 1024
V2_SUBDIR = "v2"
ASSETS_SUBDIR = "assets"


def parser_version() -> str:
    """Parser implementation stamp stored on every version row."""

    pdf_version = pptx_version = "missing"
    try:
        import pdfplumber

        pdf_version = str(getattr(pdfplumber, "__version__", "unknown"))
    except Exception:  # pragma: no cover - import guard
        log.exception("pdfplumber is not importable")
    try:
        import pptx

        pptx_version = str(getattr(pptx, "__version__", "unknown"))
    except Exception:  # pragma: no cover - import guard
        log.exception("python-pptx is not importable")
    return f"pdfplumber-{pdf_version}/python-pptx-{pptx_version}"


class DocumentError(ValueError):
    """Invalid upload or unsupported content (maps to HTTP 400)."""


class DocumentNotFound(LookupError):
    """No such document version, block set, or job (maps to HTTP 404)."""


class DocumentConflict(ValueError):
    """Same key+label already taken by different bytes (maps to HTTP 409)."""


def _text(value: Any, limit: int = 200) -> str:
    return str(value or "").strip()[:limit]


def _clean_key(value: Any, limit: int = 200) -> str:
    return str(value or "").strip()[:limit]


def _clean_applicability(value: Any) -> dict:
    if value is None:
        return {}
    if isinstance(value, str):
        try:
            value = json.loads(value) if value.strip() else {}
        except ValueError as exc:
            raise DocumentError("applicability must be a JSON object") from exc
    if not isinstance(value, dict):
        raise DocumentError("applicability must be a JSON object")
    return {str(key): value[key] for key in value}


def storage_dir(base_dir: str | Path) -> Path:
    path = Path(base_dir) / V2_SUBDIR
    path.mkdir(parents=True, exist_ok=True)
    return path


def assets_dir(base_dir: str | Path) -> Path:
    path = storage_dir(base_dir) / ASSETS_SUBDIR
    path.mkdir(parents=True, exist_ok=True)
    return path


def detect_file_type(filename: str, content: bytes) -> str:
    """Decide by magic bytes, never by extension alone."""

    head = bytes(content[:5])
    if head.startswith(PDF_MAGIC):
        return "pdf"
    if bytes(content[:4]) == PPTX_MAGIC:
        return "pptx"
    raise DocumentError(
        f"{_text(filename, 80) or 'upload'} is neither a PDF (%PDF) nor a PPTX (PK) file"
    )


def _stored_name(sha256: str, file_type: str) -> str:
    return f"{sha256}.{file_type}"


def create_version(
    conn,
    *,
    base_dir: str | Path,
    document_key: str,
    version_label: str | None,
    filename: str,
    content: bytes,
    title: str = "",
    applicability: Any = None,
    source_authenticity: str = "unverified",
    legacy_document_id: int | None = None,
) -> tuple[dict, bool]:
    """Store one immutable file version and queue its parse job.

    Returns ``(version, created)``: re-uploading identical bytes under the
    same key+label returns the stored row; different bytes under a taken
    label are a 409, never an overwrite.
    """

    key = _clean_key(document_key)
    if not key:
        raise DocumentError("document_key is required")
    label = _clean_key(version_label) or Path(str(filename or "")).stem.strip()[:200]
    if not label:
        raise DocumentError("version_label is required")
    data = bytes(content or b"")
    if not data:
        raise DocumentError("uploaded file is empty")
    if len(data) > MAX_UPLOAD_BYTES:
        raise DocumentError(
            f"uploaded file exceeds the {MAX_UPLOAD_BYTES // (1024 * 1024)} MB limit"
        )
    file_type = detect_file_type(filename, data)
    if source_authenticity not in ("unverified", "official_vendor", "confirmed_copy"):
        raise DocumentError("unknown source_authenticity")
    clean_applicability = _clean_applicability(applicability)
    digest = hashlib.sha256(data).hexdigest()

    with conn.cursor() as cur:
        cur.execute(
            """
            SELECT id, document_key, version_label, sha256, file_name,
                   file_type, file_size, stored_path, title, applicability,
                   source_authenticity, parser_version, status,
                   legacy_document_id, created_at, updated_at
            FROM v2_document_versions
            WHERE document_key=%s AND version_label=%s
            """,
            (key, label),
        )
        existing = cur.fetchone()
    if existing:
        row = dict(existing)
        if row.get("sha256") != digest:
            raise DocumentConflict(
                f"version {label!r} of {key!r} already stores different bytes; "
                "choose another version_label"
            )
        return row, False

    directory = storage_dir(base_dir)
    stored = directory / _stored_name(digest, file_type)
    if not stored.is_file():
        stored.write_bytes(data)
    with conn.cursor() as cur:
        cur.execute(
            """
            INSERT INTO v2_document_versions(
                document_key, version_label, sha256, file_name, file_type,
                file_size, stored_path, title, applicability,
                source_authenticity, parser_version, status, legacy_document_id
            ) VALUES(%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, 'uploaded', %s)
            RETURNING id, document_key, version_label, sha256, file_name,
                      file_type, file_size, stored_path, title, applicability,
                      source_authenticity, parser_version, status,
                      legacy_document_id, created_at, updated_at
            """,
            (
                key, label, digest, _text(filename, 500), file_type, len(data),
                _stored_name(digest, file_type), _text(title, 500),
                Jsonb(clean_applicability), source_authenticity,
                parser_version(), legacy_document_id,
            ),
        )
        version = dict(cur.fetchone())
    _ensure_parse_job(conn, int(version["id"]))
    return version, True


def _ensure_parse_job(conn, version_id: int) -> dict:
    with conn.cursor() as cur:
        cur.execute(
            """
            INSERT INTO v2_document_jobs(
                version_id, stage, idempotency_key, status
            ) VALUES(%s, 'parse', %s, 'queued')
            ON CONFLICT (idempotency_key) DO NOTHING
            RETURNING id, version_id, stage, block_id, checkpoint,
                      idempotency_key, status, attempts, next_run_at,
                      result_summary, error, created_at, started_at,
                      completed_at, updated_at
            """,
            (int(version_id), f"v2doc:parse:{int(version_id)}"),
        )
        row = cur.fetchone()
        if row:
            return dict(row)
        cur.execute(
            """
            SELECT id, version_id, stage, block_id, checkpoint,
                   idempotency_key, status, attempts, next_run_at,
                   result_summary, error, created_at, started_at,
                   completed_at, updated_at
            FROM v2_document_jobs WHERE idempotency_key=%s
            """,
            (f"v2doc:parse:{int(version_id)}",),
        )
        existing = cur.fetchone()
        assert existing is not None
        return dict(existing)


def get_version(conn, version_id: int) -> dict | None:
    with conn.cursor() as cur:
        cur.execute(
            """
            SELECT id, document_key, version_label, sha256, file_name,
                   file_type, file_size, stored_path, title, applicability,
                   source_authenticity, parser_version, status,
                   legacy_document_id, created_at, updated_at
            FROM v2_document_versions WHERE id=%s
            """,
            (int(version_id),),
        )
        row = cur.fetchone()
    return dict(row) if row else None


def list_versions(conn, limit: int = 100) -> list[dict]:
    with conn.cursor() as cur:
        cur.execute(
            """
            SELECT v.id, v.document_key, v.version_label, v.sha256,
                   v.file_name, v.file_type, v.file_size, v.title,
                   v.source_authenticity, v.parser_version, v.status,
                   v.created_at, v.updated_at,
                   (SELECT count(*) FROM v2_document_blocks b
                     WHERE b.version_id=v.id) AS block_count
            FROM v2_document_versions v
            ORDER BY v.id DESC
            LIMIT %s
            """,
            (max(1, min(int(limit), 500)),),
        )
        return [dict(row) for row in cur.fetchall()]


def get_blocks(conn, version_id: int) -> list[dict]:
    with conn.cursor() as cur:
        cur.execute(
            """
            SELECT b.id, b.version_id, b.block_key, b.parent_block_id,
                   b.ord, b.section_path, b.page_no, b.slide_no,
                   b.block_type, b.raw_evidence_id, b.content_hash,
                   b.layout, b.assets, b.processing_state, b.state_reason,
                   b.created_at, b.updated_at,
                   r.content AS evidence_text
            FROM v2_document_blocks b
            LEFT JOIN v2_raw_evidence r ON r.id=b.raw_evidence_id
            WHERE b.version_id=%s
            ORDER BY b.ord, b.id
            """,
            (int(version_id),),
        )
        return [dict(row) for row in cur.fetchall()]


def version_file_path(base_dir: str | Path, version: dict) -> Path:
    """Resolve the stored bytes; anything escaping storage is a 404."""

    name = str(version.get("stored_path") or "")
    candidate = (storage_dir(base_dir) / name).resolve()
    if candidate.parent != storage_dir(base_dir).resolve() or not candidate.is_file():
        raise DocumentNotFound("stored document file is unavailable")
    return candidate


# -- parsing ---------------------------------------------------------------


def _content_hash(text: str, layout: dict) -> str:
    payload = json.dumps(
        {"text": text or "", "layout": layout or {}}, ensure_ascii=False, sort_keys=True,
    )
    return hashlib.sha256(payload.encode("utf-8")).hexdigest()


def parse_file(path: str | Path, file_type: str) -> list[dict]:
    """Parse bytes on disk into block dicts (no database touched)."""

    if file_type == "pdf":
        return _parse_pdf(path)
    if file_type == "pptx":
        return _parse_pptx(path)
    raise DocumentError(f"unsupported file type: {file_type!r}")


def _parse_pdf(path: str | Path) -> list[dict]:
    import pdfplumber

    blocks: list[dict] = []
    order = 0

    def emit(**fields) -> None:
        nonlocal order
        order += 1
        block = {
            "block_key": fields.get("block_key") or f"p{fields.get('page_no')}-{order}",
            "parent_key": fields.get("parent_key"),
            "ord": order,
            "section_path": list(fields.get("section_path") or []),
            "page_no": fields.get("page_no"),
            "slide_no": None,
            "block_type": fields.get("block_type", "other"),
            "text": fields.get("text") or "",
            "layout": dict(fields.get("layout") or {}),
            "assets": list(fields.get("assets") or []),
            "needs_review_reason": fields.get("needs_review_reason") or "",
        }
        blocks.append(block)

    with pdfplumber.open(str(path)) as pdf:
        if len(pdf.pages) > MAX_PDF_PAGES:
            raise DocumentError(
                f"PDF has {len(pdf.pages)} pages, above the {MAX_PDF_PAGES} limit"
            )
        section_stack: list[str] = []
        for page_number, page in enumerate(pdf.pages, start=1):
            lines = _pdf_text_lines(page)
            tables = page.extract_tables() or []
            images = list(getattr(page, "images", None) or [])
            body_size = _pdf_body_size(lines)
            page_has_text = any(line["text"] for line in lines)
            for line in lines:
                if not line["text"]:
                    continue
                if _pdf_is_heading(line, body_size):
                    section_stack = [*section_stack[: _pdf_heading_level(line, body_size)], line["text"][:200]]
                    emit(page_no=page_number, block_type="heading",
                         text=line["text"], section_path=list(section_stack),
                         layout={"font_size": line["size"]})
                else:
                    emit(page_no=page_number, block_type="paragraph",
                         text=line["text"], section_path=list(section_stack),
                         layout={"font_size": line["size"]})
            for table_index, table in enumerate(tables, start=1):
                cells = [
                    {"row": r, "col": c, "row_span": 1, "col_span": 1,
                     "text": str(cell or "").strip()}
                    for r, row in enumerate(table)
                    for c, cell in enumerate(row)
                ]
                if not any(cell["text"] for cell in cells):
                    continue
                text = "\n".join(
                    " | ".join(cell["text"] for cell in cells if cell["row"] == r)
                    for r in range(len(table))
                ).strip()
                emit(page_no=page_number,
                     block_key=f"p{page_number}-table-{table_index}",
                     block_type="table", text=text,
                     section_path=list(section_stack),
                     layout={"rows": len(table),
                             "columns": max((len(row) for row in table), default=0),
                             "cells": cells,
                             "header_row": 0})
            for image_index, image in enumerate(images, start=1):
                emit(page_no=page_number,
                     block_key=f"p{page_number}-img-{image_index}",
                     block_type="image", text="",
                     section_path=list(section_stack),
                     layout={},
                     assets=[{
                         "kind": "image",
                         "asset_ref": f"page-{page_number}-img-{image_index}",
                         "bbox": [image.get("x0"), image.get("top"),
                                  image.get("x1"), image.get("bottom")],
                         "interpretation_status": "unexplained",
                     }],
                     needs_review_reason=(
                         "image_only_page" if not page_has_text and not tables else ""
                     ))
            if not page_has_text and not tables and not images:
                emit(page_no=page_number, block_type="other", text="",
                     needs_review_reason="no_extractable_content")
    return blocks


def _pdf_text_lines(page) -> list[dict]:
    """Group page chars into lines with text and dominant font size."""

    buckets: dict[int, list] = {}
    for char in getattr(page, "chars", None) or []:
        text = char.get("text") or ""
        # Keep inter-word spaces for joining; blank positioning artifacts
        # (empty text) carry no content.
        if text:
            buckets.setdefault(round(float(char.get("top", 0))), []).append(char)
    lines = []
    for top in sorted(buckets):
        chars = sorted(buckets[top], key=lambda item: float(item.get("x0", 0)))
        text = "".join(item.get("text") or "" for item in chars).strip()
        content = [item for item in chars if (item.get("text") or "").strip()]
        sizes = sorted(float(item.get("size", 0) or 0) for item in content)
        size = sizes[len(sizes) // 2] if sizes else 0.0
        lines.append({"text": text, "size": size})
    return lines


def _pdf_body_size(lines: list[dict]) -> float:
    weighted: dict[float, int] = {}
    for line in lines:
        if line["text"]:
            weighted[line["size"]] = weighted.get(line["size"], 0) + len(line["text"])
    if not weighted:
        return 0.0
    return max(weighted, key=lambda size: weighted[size])


def _pdf_is_heading(line: dict, body_size: float) -> bool:
    text = line["text"]
    return bool(
        body_size
        and line["size"] >= body_size + 1.0
        and 0 < len(text) <= 200
    )


def _pdf_heading_level(line: dict, body_size: float) -> int:
    # Larger headings outrank smaller ones in the section stack; equal-size
    # headings are siblings replacing each other.
    return 0 if line["size"] > body_size + 4.0 else 1


def _parse_pptx(path: str | Path) -> list[dict]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    slides = list(presentation.slides)
    if len(slides) > MAX_PPTX_SLIDES:
        raise DocumentError(
            f"PPTX has {len(slides)} slides, above the {MAX_PPTX_SLIDES} limit"
        )
    blocks: list[dict] = []
    order = 0

    def emit(**fields) -> None:
        nonlocal order
        order += 1
        block = {
            "block_key": fields.get("block_key") or f"s{fields.get('slide_no')}-{order}",
            "parent_key": fields.get("parent_key"),
            "ord": order,
            "section_path": list(fields.get("section_path") or []),
            "page_no": None,
            "slide_no": fields.get("slide_no"),
            "slide_title": fields.get("slide_title") or "",
            "block_type": fields.get("block_type", "other"),
            "text": fields.get("text") or "",
            "layout": dict(fields.get("layout") or {}),
            "assets": list(fields.get("assets") or []),
            "needs_review_reason": fields.get("needs_review_reason") or "",
        }
        blocks.append(block)

    for slide_number, slide in enumerate(slides, start=1):
        title = ""
        try:
            if slide.shapes.title and (slide.shapes.title.text or "").strip():
                title = slide.shapes.title.text.strip()
        except Exception:
            title = ""
        section = [title[:200]] if title else [f"Slide {slide_number}"]
        notes = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes = ""
        shapes = sorted(
            (shape for shape in slide.shapes),
            key=lambda item: (int(getattr(item, "top", 0) or 0),
                              int(getattr(item, "left", 0) or 0)),
        )
        emitted = 0
        image_index = 0
        for position, shape in enumerate(shapes):
            base_layout = {"slide_title": title, "z_order": position}
            try:
                if shape.has_table:
                    cells, skipped = _pptx_table_cells(shape.table)
                    if not any(cell["text"] for cell in cells):
                        continue
                    text = "\n".join(
                        " | ".join(
                            cell["text"] for cell in cells
                            if cell["row"] == row
                        )
                        for row in sorted({cell["row"] for cell in cells})
                    ).strip()
                    emit(slide_no=slide_number, slide_title=title,
                         block_type="table", text=text, section_path=section,
                         layout={**base_layout,
                                 "rows": len(shape.table.rows),
                                 "columns": len(shape.table.columns),
                                 "cells": cells, "header_row": 0,
                                 "merged_skipped": skipped})
                    emitted += 1
                elif getattr(shape, "image", None) is not None:
                    image_index += 1
                    emit(slide_no=slide_number, slide_title=title,
                         block_type="image", text="", section_path=section,
                         layout=base_layout,
                         assets=[{
                             "kind": "image",
                             "asset_ref": f"slide-{slide_number}-img-{image_index}",
                             "image_bytes": shape.image.blob,
                             "image_ext": str(getattr(shape.image, "ext", "png") or "png"),
                             "interpretation_status": "unexplained",
                         }])
                    emitted += 1
                elif bool(getattr(shape, "has_chart", False)):
                    emit(slide_no=slide_number, slide_title=title,
                         block_type="other", text="", section_path=section,
                         layout=base_layout,
                         assets=[{"kind": "chart",
                                  "interpretation_status": "unexplained"}],
                         needs_review_reason="unexplained_graphic")
                    emitted += 1
                elif bool(getattr(shape, "has_text_frame", False)):
                    text = (shape.text_frame.text or "").strip()
                    if not text or (title and text == title and position == 0):
                        if title and text == title:
                            emit(slide_no=slide_number, slide_title=title,
                                 block_key=f"s{slide_number}-title",
                                 block_type="heading", text=title,
                                 section_path=section, layout=base_layout)
                            emitted += 1
                        continue
                    emit(slide_no=slide_number, slide_title=title,
                         block_type="paragraph", text=text,
                         section_path=section, layout=base_layout)
                    emitted += 1
                else:
                    emit(slide_no=slide_number, slide_title=title,
                         block_type="other", text="", section_path=section,
                         layout={**base_layout,
                                 "shape_type": str(getattr(shape, "shape_type", ""))},
                         assets=[{"kind": "shape",
                                  "interpretation_status": "unexplained"}],
                         needs_review_reason="unexplained_graphic")
                    emitted += 1
            except DocumentError:
                raise
            except Exception as exc:  # one bad shape never kills the slide
                log.warning("unreadable shape sidelined slide=%s: %s", slide_number, exc)
                emit(slide_no=slide_number, slide_title=title,
                     block_type="other", text="", section_path=section,
                     layout=base_layout,
                     needs_review_reason="unreadable_shape")
                emitted += 1
        if notes:
            emit(slide_no=slide_number, slide_title=title,
                 block_key=f"s{slide_number}-notes",
                 block_type="note", text=notes, section_path=section,
                 layout={"slide_title": title, "speaker_notes": True})
            emitted += 1
        if not emitted:
            emit(slide_no=slide_number, slide_title=title,
                 block_type="other", text="", section_path=section,
                 needs_review_reason="empty_slide")
    return blocks


def _pptx_table_cells(table) -> tuple[list[dict], int]:
    """Cells with merge spans; spanned (non-origin) cells are counted out."""

    cells: list[dict] = []
    skipped = 0
    for row_index, row in enumerate(table.rows):
        for col_index, cell in enumerate(row.cells):
            try:
                spanned = bool(cell.is_spanned)
            except Exception:
                spanned = False
            if spanned:
                skipped += 1
                continue
            try:
                row_span = int(cell.span_height or 1)
            except Exception:
                row_span = 1
            try:
                col_span = int(cell.span_width or 1)
            except Exception:
                col_span = 1
            cells.append({
                "row": row_index, "col": col_index,
                "row_span": max(1, row_span), "col_span": max(1, col_span),
                "text": (cell.text or "").strip(),
            })
    return cells, skipped


# -- persistence -----------------------------------------------------------


def save_parsed_blocks(
    conn, version_id: int, blocks: list[dict], *, base_dir: str | Path,
) -> dict:
    """Persist blocks with raw-evidence text; returns coverage counts."""

    version = get_version(conn, version_id)
    if version is None:
        raise DocumentNotFound(f"V2 document version {int(version_id)} was not found")
    locator = f"v2-doc:{int(version_id)}"
    counts = {"blocks": 0, "evidence": 0, "assets": 0, "needs_review": 0}
    asset_bytes = 0
    with conn.cursor() as cur:
        for block in blocks:
            text = str(block.get("text") or "")
            page_no = block.get("page_no")
            slide_no = block.get("slide_no")
            if slide_no is not None:
                position = f"slide {int(slide_no)}"
            elif page_no is not None:
                position = f"page {int(page_no)}"
            else:
                position = "unknown position"
            evidence_id = None
            if text:
                cur.execute(
                    """
                    INSERT INTO v2_raw_evidence(
                        evidence_type, author_role, content, raw_payload,
                        source_label, source_locator
                    ) VALUES('document', 'unknown', %s, %s, %s, %s)
                    RETURNING id
                    """,
                    (
                        text,
                        Jsonb({
                            "document_version_id": int(version_id),
                            "document_key": str(version.get("document_key") or ""),
                            "block_key": str(block.get("block_key") or ""),
                            "page_no": page_no, "slide_no": slide_no,
                        }),
                        _text(version.get("title") or version.get("file_name"), 500),
                        f"{locator}:{position}",
                    ),
                )
                evidence_id = int(cur.fetchone()["id"])
                counts["evidence"] += 1
            assets = []
            for asset in block.get("assets") or []:
                item = dict(asset)
                blob = item.pop("image_bytes", None)
                if blob:
                    if asset_bytes + len(blob) > MAX_ASSET_BYTES:
                        raise DocumentError("document image assets exceed the size budget")
                    digest = hashlib.sha256(bytes(blob)).hexdigest()
                    ext = "".join(
                        ch for ch in str(item.get("image_ext") or "png").lower()
                        if ch.isalnum()
                    )[:10] or "png"
                    target = assets_dir(base_dir) / f"{digest}.{ext}"
                    if not target.is_file():
                        target.write_bytes(bytes(blob))
                    asset_bytes += len(blob)
                    item["asset_ref"] = f"assets/{digest}.{ext}"
                    item["blob_sha256"] = digest
                    counts["assets"] += 1
                assets.append(item)
            reason = str(block.get("needs_review_reason") or "")
            if reason:
                counts["needs_review"] += 1
            cur.execute(
                """
                INSERT INTO v2_document_blocks(
                    version_id, block_key, ord, section_path, page_no,
                    slide_no, block_type, raw_evidence_id, content_hash,
                    layout, assets, processing_state, state_reason
                ) VALUES(%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
                ON CONFLICT (version_id, block_key) DO NOTHING
                """,
                (
                    int(version_id), _text(block.get("block_key"), 200),
                    int(block.get("ord") or 0),
                    [str(part)[:200] for part in (block.get("section_path") or [])],
                    page_no, slide_no, str(block.get("block_type") or "other"),
                    evidence_id,
                    _content_hash(text, block.get("layout") or {}),
                    Jsonb(block.get("layout") or {}), Jsonb(assets),
                    "needs_review" if reason else "pending", reason[:1000],
                ),
            )
            counts["blocks"] += 1
        cur.execute(
            """
            UPDATE v2_document_versions
            SET status='parsed', updated_at=CURRENT_TIMESTAMP
            WHERE id=%s
            """,
            (int(version_id),),
        )
    return counts


def fail_version(conn, version_id: int, reason: str) -> None:
    with conn.cursor() as cur:
        cur.execute(
            """
            UPDATE v2_document_versions
            SET status='parse_failed', updated_at=CURRENT_TIMESTAMP
            WHERE id=%s
            """,
            (int(version_id),),
        )


# -- jobs ------------------------------------------------------------------


_JOB_COLUMNS = (
    "id, version_id, stage, block_id, checkpoint, idempotency_key, "
    "status, attempts, next_run_at, result_summary, error, created_at, "
    "started_at, completed_at, updated_at"
)

def get_document_job(conn, job_id: int) -> dict | None:
    with conn.cursor() as cur:
        cur.execute(
            f"SELECT {_JOB_COLUMNS} FROM v2_document_jobs WHERE id=%s",
            (int(job_id),),
        )
        row = cur.fetchone()
    return dict(row) if row else None


def latest_document_job(conn, version_id: int, stage: str = "parse") -> dict | None:
    """Newest job for one version+stage, for detail views and retries."""

    with conn.cursor() as cur:
        cur.execute(
            f"""
            SELECT {_JOB_COLUMNS} FROM v2_document_jobs
            WHERE version_id=%s AND stage=%s
            ORDER BY id DESC LIMIT 1
            """,
            (int(version_id), str(stage)),
        )
        row = cur.fetchone()
    return dict(row) if row else None


def version_coverage(conn, version_id: int) -> dict:
    """Whole-document disposition: every block states its destination.

    ``complete`` means every block has a destination (knowledge, proposal,
    evidence_only, needs_review) and no learn job is still open -- never a
    claim about understanding quality.
    """

    with conn.cursor() as cur:
        cur.execute(
            """
            SELECT processing_state, count(*) AS n
            FROM v2_document_blocks
            WHERE version_id=%s
            GROUP BY processing_state
            """,
            (int(version_id),),
        )
        states = {str(row["processing_state"]): int(row["n"]) for row in cur.fetchall()}
        cur.execute(
            """
            SELECT id, block_key, page_no, slide_no, block_type,
                   processing_state, state_reason
            FROM v2_document_blocks
            WHERE version_id=%s AND processing_state IN ('pending', 'parse_failed', 'learn_failed')
            ORDER BY ord, id
            """,
            (int(version_id),),
        )
        unfinished = [dict(row) for row in cur.fetchall()]
        cur.execute(
            """
            SELECT stage, status, count(*) AS n
            FROM v2_document_jobs
            WHERE version_id=%s
            GROUP BY stage, status
            """,
            (int(version_id),),
        )
        jobs = [
            {"stage": str(row["stage"]), "status": str(row["status"]), "count": int(row["n"])}
            for row in cur.fetchall()
        ]
    total = sum(states.values())
    with_destination = total - states.get("pending", 0)
    return {
        "version_id": int(version_id),
        "total_blocks": total,
        "by_state": states,
        "with_destination": with_destination,
        "unfinished_blocks": unfinished,
        "jobs": jobs,
        "complete": not unfinished and not any(
            item["stage"] == "learn" and item["status"] in ("queued", "processing")
            for item in jobs
        ),
    }


def claim_document_job(conn, stages: tuple[str, ...] = ("parse",)) -> dict | None:
    """Claim one due job; the worker holds no lock across parsing."""

    with conn.cursor() as cur:
        cur.execute(
            f"""
            SELECT {_JOB_COLUMNS} FROM v2_document_jobs
            WHERE status='queued' AND next_run_at <= CURRENT_TIMESTAMP
              AND stage = ANY(%s)
            ORDER BY next_run_at, id LIMIT 1
            FOR UPDATE SKIP LOCKED
            """,
            (list(stages),),
        )
        row = cur.fetchone()
        if not row:
            return None
        job = dict(row)
        cur.execute(
            """
            UPDATE v2_document_jobs
            SET status='processing', attempts=attempts+1,
                started_at=CURRENT_TIMESTAMP, updated_at=CURRENT_TIMESTAMP
            WHERE id=%s
            """,
            (int(job["id"]),),
        )
    result = get_document_job(conn, int(job["id"]))
    assert result is not None
    return result


def complete_document_job(conn, job_id: int, summary: dict) -> None:
    with conn.cursor() as cur:
        cur.execute(
            """
            UPDATE v2_document_jobs
            SET status='completed', result_summary=%s, error='',
                completed_at=CURRENT_TIMESTAMP, updated_at=CURRENT_TIMESTAMP
            WHERE id=%s
            """,
            (Jsonb(dict(summary or {})), int(job_id)),
        )


def fail_document_job(conn, job_id: int, error: str, *, retryable: bool) -> None:
    """Retryable failures come back via next_run_at; terminal ones stay failed."""

    with conn.cursor() as cur:
        if retryable:
            cur.execute(
                """
                UPDATE v2_document_jobs
                SET status='queued', error=%s,
                    next_run_at=CURRENT_TIMESTAMP + make_interval(secs => attempts * 60),
                    updated_at=CURRENT_TIMESTAMP
                WHERE id=%s
                """,
                (str(error or "")[:2000], int(job_id)),
            )
        else:
            cur.execute(
                """
                UPDATE v2_document_jobs
                SET status='failed', error=%s,
                    completed_at=CURRENT_TIMESTAMP, updated_at=CURRENT_TIMESTAMP
                WHERE id=%s
                """,
                (str(error or "")[:2000], int(job_id)),
            )


def retry_document_job(conn, job_id: int) -> dict | None:
    with conn.cursor() as cur:
        cur.execute(
            """
            UPDATE v2_document_jobs
            SET status='queued', error='', next_run_at=CURRENT_TIMESTAMP,
                updated_at=CURRENT_TIMESTAMP
            WHERE id=%s AND status='failed'
            RETURNING id
            """,
            (int(job_id),),
        )
        updated = cur.fetchone()
    if not updated:
        return get_document_job(conn, job_id)
    result = get_document_job(conn, int(updated["id"]))
    assert result is not None
    return result


def unfinished_document_job_ids(conn) -> list[int]:
    with conn.cursor() as cur:
        cur.execute(
            "SELECT id FROM v2_document_jobs WHERE status='processing' ORDER BY id"
        )
        return [int(row["id"]) for row in cur.fetchall()]


__all__ = [
    "ASSETS_SUBDIR",
    "FILE_TYPES",
    "MAX_UPLOAD_BYTES",
    "V2_SUBDIR",
    "DocumentConflict",
    "DocumentError",
    "DocumentNotFound",
    "assets_dir",
    "claim_document_job",
    "complete_document_job",
    "create_version",
    "detect_file_type",
    "fail_document_job",
    "fail_version",
    "get_blocks",
    "get_document_job",
    "get_version",
    "latest_document_job",
    "list_versions",
    "parse_file",
    "parser_version",
    "retry_document_job",
    "save_parsed_blocks",
    "storage_dir",
    "affected_knowledge",
    "compare_versions",
    "lineage_version_ids",
    "record_revalidation",
    "unfinished_document_job_ids",
    "version_coverage",
    "version_file_path",
]


# -- version comparison + revalidation (Phase 5.2) ----------------------------


GLOBAL_SECTION_KEYWORDS = (
    "警告", "注意", "适用", "范围", "前置", "限制",
    "warning", "caution", "scope", "applicability", "prerequisite",
    "limitation", "requirement",
)


def _section_text_signature(text: str) -> str:
    return hashlib.sha256(_normalized_section_text(text).encode("utf-8")).hexdigest()


def _normalized_section_text(text: str) -> str:
    import re as _re

    collapsed = _re.sub(r"\s+", " ", str(text or "").casefold()).strip()
    return collapsed


def _is_global_section(title: str) -> bool:
    folded = str(title or "").casefold()
    return any(keyword in folded for keyword in GLOBAL_SECTION_KEYWORDS)


def compare_versions(conn, old_version_id: int, new_version_id: int) -> dict:
    """Section-level comparison between two parsed versions.

    Blocks group by their top section (or slide title); page/slide numbers
    never count as changes.  Returns added/changed/removed/unmatched section
    lists with block references, plus whether a global scope/warning section
    changed (which conservatively affects every dependent unit).
    """

    old_blocks = get_blocks(conn, old_version_id)
    new_blocks = get_blocks(conn, new_version_id)
    if not old_blocks:
        raise DocumentNotFound(f"V2 document version {int(old_version_id)} has no blocks")
    if not new_blocks:
        raise DocumentNotFound(f"V2 document version {int(new_version_id)} has no blocks")

    def sections(blocks: list[dict]) -> dict[str, list[dict]]:
        # Full heading path: a changed subsection must not lump the whole
        # chapter, and page/slide numbers never participate.
        grouped: dict[str, list[dict]] = {}
        for block in blocks:
            path = list(block.get("section_path") or [])
            title = " / ".join(str(part) for part in path if str(part).strip())
            grouped.setdefault(title or "Untitled", []).append(block)
        return grouped

    old_sections = sections(old_blocks)
    new_sections = sections(new_blocks)
    added, changed, removed, unmatched = [], [], [], []
    for title, blocks in sorted(new_sections.items()):
        if title not in old_sections:
            added.append(_section_view(title, blocks))
            continue
        old_sig = _section_text_signature(
            "\n".join(str(block.get("evidence_text") or "") for block in old_sections[title]))
        new_sig = _section_text_signature(
            "\n".join(str(block.get("evidence_text") or "") for block in blocks))
        if old_sig != new_sig:
            changed.append(_section_view(title, blocks))
    for title, blocks in sorted(old_sections.items()):
        if title not in new_sections:
            removed.append(_section_view(title, blocks))
        elif not any(
            str(block.get("evidence_text") or "").strip()
            for block in new_sections[title]
        ) and any(
            str(block.get("evidence_text") or "").strip() for block in blocks
        ):
            unmatched.append(_section_view(title, blocks))
    global_changed = any(
        _is_global_section(item["section"])
        for item in changed
    )
    return {
        "old_version_id": int(old_version_id),
        "new_version_id": int(new_version_id),
        "added": added,
        "changed": changed,
        "removed": removed,
        "unmatched": unmatched,
        "global_scope_changed": global_changed,
    }


def _section_view(title: str, blocks: list[dict]) -> dict:
    return {
        "section": title,
        "global": _is_global_section(title),
        "blocks": [
            {
                "block_id": int(block["id"]),
                "block_key": str(block.get("block_key") or ""),
                "page_no": block.get("page_no"),
                "slide_no": block.get("slide_no"),
                "block_type": str(block.get("block_type") or ""),
            }
            for block in blocks
        ],
    }


def affected_knowledge(conn, old_version_id: int, comparison: dict) -> list[dict]:
    """Knowledge rows of the old version whose evidence sits in changed areas.

    A changed global scope/warning section conservatively affects every unit
    derived from the old version.  Otherwise only units citing blocks of
    added/changed/removed/unmatched sections are listed.  Read-only: nothing
    is deactivated here.
    """

    sections = {
        item["section"]
        for group in ("changed", "removed", "unmatched")
        for item in comparison.get(group, [])
    }
    with conn.cursor() as cur:
        if comparison.get("global_scope_changed"):
            cur.execute(
                """
                SELECT k.id, k.title, k.trust, k.validation_status, k.revision
                FROM v2_knowledge k
                WHERE k.origin_document_version_id=%s AND k.active=TRUE
                ORDER BY k.id
                """,
                (int(old_version_id),),
            )
            return [dict(row) for row in cur.fetchall()]
        if not sections:
            return []
        cur.execute(
            """
            SELECT DISTINCT k.id, k.title, k.trust, k.validation_status, k.revision
            FROM v2_knowledge k
            JOIN v2_knowledge_sources s ON s.knowledge_id=k.id
            JOIN v2_raw_evidence r ON r.id=s.raw_evidence_id
            JOIN v2_document_blocks b ON b.raw_evidence_id=r.id
            WHERE k.origin_document_version_id=%s
              AND k.active=TRUE
              AND s.active=TRUE AND s.relation='supports'
              AND b.version_id=%s
              AND array_to_string(b.section_path, ' / ') = ANY(%s)
            ORDER BY k.id
            """,
            (int(old_version_id), int(old_version_id), sorted(sections)),
        )
        return [dict(row) for row in cur.fetchall()]


def lineage_version_ids(conn, version_id: int) -> list[int]:
    """The version plus its confirmed predecessors, newest first."""

    chain = []
    seen = set()
    current: int | None = int(version_id)
    with conn.cursor() as cur:
        while current is not None and current not in seen:
            seen.add(current)
            chain.append(current)
            cur.execute(
                "SELECT previous_version_id FROM v2_document_versions WHERE id=%s",
                (current,),
            )
            row = cur.fetchone()
            current = int(row["previous_version_id"]) if row and row.get("previous_version_id") else None
    return chain


def record_revalidation(conn, new_version_id: int, previous_version_id: int,
                         comparison: dict, affected: list[dict]) -> dict:
    """Link the new version to its predecessor and store the comparison.

    Old units are never touched here: they keep serving until their own
    evidence fails or an engineer explicitly revalidates them.
    """

    if int(new_version_id) == int(previous_version_id):
        raise DocumentError("a version cannot revalidate itself")
    with conn.cursor() as cur:
        cur.execute(
            "SELECT id, document_key FROM v2_document_versions WHERE id=%s",
            (int(previous_version_id),),
        )
        previous = cur.fetchone()
        cur.execute(
            "SELECT id, document_key FROM v2_document_versions WHERE id=%s",
            (int(new_version_id),),
        )
        current = cur.fetchone()
    if not previous or not current:
        raise DocumentNotFound("both versions must exist for revalidation")
    if str(previous["document_key"] or "") != str(current["document_key"] or ""):
        raise DocumentError("revalidation links versions of the same document_key only")
    summary = {
        "previous_version_id": int(previous_version_id),
        "added": [item["section"] for item in comparison.get("added", [])],
        "changed": [item["section"] for item in comparison.get("changed", [])],
        "removed": [item["section"] for item in comparison.get("removed", [])],
        "unmatched": [item["section"] for item in comparison.get("unmatched", [])],
        "global_scope_changed": bool(comparison.get("global_scope_changed")),
        "affected_knowledge_ids": [int(item["id"]) for item in affected],
    }
    with conn.cursor() as cur:
        cur.execute(
            """
            UPDATE v2_document_versions
            SET previous_version_id=%s, change_summary=%s,
                updated_at=CURRENT_TIMESTAMP
            WHERE id=%s
            RETURNING id, document_key, version_label, previous_version_id,
                      change_summary, status
            """,
            (int(previous_version_id), Jsonb(summary), int(new_version_id)),
        )
        row = cur.fetchone()
    assert row is not None
    return dict(row)
